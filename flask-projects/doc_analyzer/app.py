import os
import csv

from flask import Flask, render_template, request, jsonify
from werkzeug.utils import secure_filename

# ── text-extraction libraries ──────────────────────────────────────────────────
import pdfplumber                    # PDF
from docx import Document as DocxDoc # DOCX
from pptx import Presentation        # PPTX
import openpyxl                      # XLSX
from PIL import Image                # image open + EXIF
import pytesseract                   # OCR wrapper
from PIL.ExifTags import TAGS        # human-readable EXIF tag names
# ──────────────────────────────────────────────────────────────────────────────

# ── auto-detect Tesseract on Windows ──────────────────────────────────────────
_WIN_TESS_PATHS = [
    r"C:\Program Files\Tesseract-OCR\tesseract.exe",
    r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
]
for _p in _WIN_TESS_PATHS:
    if os.path.isfile(_p):
        pytesseract.pytesseract.tesseract_cmd = _p
        break
# ──────────────────────────────────────────────────────────────────────────────

app = Flask(__name__)

app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024   # 50 MB upload limit
app.config['UPLOAD_FOLDER'] = os.path.join(os.path.dirname(__file__), 'uploads')

ALLOWED_EXTENSIONS = {
    'pdf', 'docx', 'doc',
    'pptx', 'ppt',
    'xlsx', 'xls',
    'csv',
    'txt', 'md', 'log', 'rtf',
    'py', 'js', 'ts', 'html', 'css',
    'json', 'xml', 'yaml', 'yml',
    'ini', 'cfg', 'env',
    # images (OCR)
    'png', 'jpg', 'jpeg', 'bmp', 'tiff', 'tif', 'gif', 'webp',
}

IMAGE_EXTENSIONS = {'png', 'jpg', 'jpeg', 'bmp', 'tiff', 'tif', 'gif', 'webp'}

os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)


# ── helpers ────────────────────────────────────────────────────────────────────

def allowed_file(filename: str) -> bool:
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS


def _safe_str(v) -> str:
    """Convert any metadata value to a JSON-safe string."""
    try:
        return str(v)
    except Exception:
        return ''


def extract_text(filepath: str, filename: str) -> dict:
    ext = filename.rsplit('.', 1)[1].lower()

    result = {
        'content':   '',
        'pages':     0,
        'words':     0,
        'chars':     0,
        'file_type': ext.upper(),
        'metadata':  {},
        'structure': [],
        'error':     None,
    }

    try:
        # ── PDF ────────────────────────────────────────────────────────────────
        if ext == 'pdf':
            with pdfplumber.open(filepath) as pdf:
                result['pages'] = len(pdf.pages)

                # clean metadata
                raw_meta = pdf.metadata or {}
                result['metadata'] = {k: _safe_str(v) for k, v in raw_meta.items() if v}

                parts = []
                for i, page in enumerate(pdf.pages, 1):
                    text = page.extract_text() or ''
                    parts.append(f"━━━  Page {i}  ━━━\n{text}")
                    result['structure'].append({'type': 'page', 'number': i, 'chars': len(text)})

                result['content'] = '\n\n'.join(parts)

        # ── DOCX ───────────────────────────────────────────────────────────────
        elif ext in ('docx', 'doc'):
            doc = DocxDoc(filepath)
            parts = []

            for para in doc.paragraphs:
                if para.text.strip():
                    parts.append(para.text)
                    result['structure'].append({'type': 'paragraph', 'style': para.style.name})

            for i, table in enumerate(doc.tables, 1):
                rows = []
                for row in table.rows:
                    rows.append(' | '.join(c.text.strip() for c in row.cells))
                parts.append(f"\n[Table {i}]\n" + '\n'.join(rows))

            result['content'] = '\n'.join(parts)
            result['pages']   = len(doc.sections)

            # core properties
            cp = doc.core_properties
            result['metadata'] = {
                k: _safe_str(v) for k, v in {
                    'Author':   cp.author,
                    'Title':    cp.title,
                    'Subject':  cp.subject,
                    'Created':  cp.created,
                    'Modified': cp.modified,
                    'Keywords': cp.keywords,
                }.items() if v
            }

        # ── PPTX ───────────────────────────────────────────────────────────────
        elif ext in ('pptx', 'ppt'):
            prs = Presentation(filepath)
            result['pages'] = len(prs.slides)
            parts = []

            for i, slide in enumerate(prs.slides, 1):
                slide_text = f"━━━  Slide {i}  ━━━\n"
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        slide_text += shape.text.strip() + '\n'
                parts.append(slide_text)
                result['structure'].append({'type': 'slide', 'number': i})

            result['content'] = '\n\n'.join(parts)

        # ── XLSX ───────────────────────────────────────────────────────────────
        elif ext in ('xlsx', 'xls'):
            wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
            parts = []

            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                lines = [f"━━━  Sheet: {sheet_name}  ━━━"]
                for row in ws.iter_rows(values_only=True):
                    row_str = ' | '.join(str(c) if c is not None else '' for c in row)
                    if row_str.replace('|', '').strip():
                        lines.append(row_str)
                parts.append('\n'.join(lines))
                result['structure'].append({'type': 'sheet', 'name': sheet_name})

            result['content'] = '\n\n'.join(parts)
            result['pages']   = len(wb.sheetnames)

        # ── CSV ────────────────────────────────────────────────────────────────
        elif ext == 'csv':
            with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
                raw = f.read()

            reader = csv.reader(raw.splitlines())
            rows   = list(reader)
            result['pages']   = len(rows)
            result['content'] = '\n'.join(' | '.join(row) for row in rows)

        # ── Images  (OCR via Tesseract) ────────────────────────────────────────
        elif ext in IMAGE_EXTENSIONS:
            img = Image.open(filepath)

            # ── image metadata ────────────────────────────────────────────────
            meta = {
                'Format': img.format or ext.upper(),
                'Mode':   img.mode,
                'Width':  f"{img.size[0]} px",
                'Height': f"{img.size[1]} px",
            }

            # EXIF (JPEG / TIFF)
            try:
                exif_raw = img._getexif()  # returns None for PNG/BMP etc.
                if exif_raw:
                    for tag_id, value in exif_raw.items():
                        tag = TAGS.get(tag_id, tag_id)
                        if isinstance(value, (str, int, float, bytes)):
                            if str(tag) not in ('MakerNote', 'UserComment', 'PrintImageMatching'):
                                meta[str(tag)] = str(value)[:200]
            except Exception:
                pass

            result['metadata'] = meta
            result['pages']    = 1

            # ── OCR ───────────────────────────────────────────────────────────
            try:
                ocr_text = pytesseract.image_to_string(img, lang='eng')
                result['content'] = ocr_text.strip() if ocr_text.strip() else \
                    "[No text detected in this image]"
            except pytesseract.TesseractNotFoundError:
                result['content'] = (
                    "⚠️  Tesseract OCR engine is not installed on your system.\n\n"
                    "To enable image text extraction, install Tesseract:\n\n"
                    "  Windows → https://github.com/UB-Mannheim/tesseract/wiki\n"
                    "            (Download the installer, run it, restart the server)\n\n"
                    "  Linux   → sudo apt install tesseract-ocr\n"
                    "  macOS   → brew install tesseract\n\n"
                    "Image info (width, height, format) is shown in the Metadata panel below."
                )
                result['error'] = 'Tesseract not installed'

        # ── Plain text / code / config ─────────────────────────────────────────
        else:
            with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
                result['content'] = f.read()
            result['pages'] = result['content'].count('\n') + 1

    except Exception as exc:
        result['content'] = f"[Extraction error]\n{exc}"
        result['error']   = str(exc)

    # ── common stats ──────────────────────────────────────────────────────────
    text = result['content']
    result['chars'] = len(text)
    result['words'] = len(text.split()) if text.strip() else 0

    return result


# ── routes ─────────────────────────────────────────────────────────────────────

@app.route('/')
def index():
    return render_template('index.html')


@app.route('/upload', methods=['POST'])
def upload():
    if 'file' not in request.files:
        return jsonify({'error': 'No file attached to the request.'}), 400

    file = request.files['file']

    if not file.filename:
        return jsonify({'error': 'No file selected.'}), 400

    if not allowed_file(file.filename):
        return jsonify({
            'error': f'Unsupported file type. Supported: {", ".join(sorted(ALLOWED_EXTENSIONS))}'
        }), 400

    filename = secure_filename(file.filename)
    filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    file.save(filepath)

    result             = extract_text(filepath, file.filename)
    result['filename'] = file.filename
    result['size']     = os.path.getsize(filepath)

    os.remove(filepath)   # clean up — we only needed it for parsing

    return jsonify(result)


# ── entry point ────────────────────────────────────────────────────────────────

if __name__ == '__main__':
    app.run(debug=True, port=5001)
